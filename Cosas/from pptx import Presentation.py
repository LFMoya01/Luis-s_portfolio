from pptx import Presentation
from pptx.util import Inches

# Crear la presentación
prs = Presentation()

# Título de la presentación
slide_layout = prs.slide_layouts[0] # Diseño de título
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Derechos Humanos: Historia, Interpretaciones y Desafíos"
subtitle.text = "Una presentación académica"

# Slide 1: Definición de Derechos Humanos
slide_layout = prs.slide_layouts[1] # Diseño de título y contenido
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Definición de Derechos Humanos"
content.text = ("Los derechos humanos son un conjunto de normas morales que aseguran un trato mínimo digno "
                "a cualquier individuo, independientemente de sus cualidades físicas, sociales, opiniones políticas, entre otros.")

# Slide 2: Importancia de los Derechos Humanos
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Importancia de los Derechos Humanos"
content.text = ("Los derechos humanos son fundamentales para la convivencia pacífica y el respeto mutuo entre los seres humanos.")

# Slide 3: Historia y Evolución de los Derechos Humanos
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Historia y Evolución de los Derechos Humanos"
content.text = ("- Antigüedad: Código de Hammurabi, la Carta Magna.\n"
                "- Edad Moderna: Declaración de Independencia de los EE.UU., Declaración de los Derechos del Hombre y del Ciudadano en Francia.\n"
                "- Edad Contemporánea: Declaración Universal de los Derechos Humanos (1948).")

# Slide 4: Interpretaciones y Violaciones de los Derechos Humanos
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Interpretaciones y Violaciones de los Derechos Humanos"
content.text = ("- Interpretaciones: Diferencias en la interpretación de los derechos humanos en distintas culturas y sistemas políticos.\n"
                "- Violaciones Históricas y Actuales: Segregación racial, esclavitud, genocidio, violencia de género, discriminación, conflictos armados.")

# Slide 5: Origen Multivariable de las Violaciones de Derechos Humanos
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Origen Multivariable de las Violaciones de Derechos Humanos"
content.text = ("- Segregación: Apartheid en Sudáfrica, segregación en Estados Unidos.\n"
                "- Esclavitud: Historia de la esclavitud, trata transatlántica de esclavos, esclavitud moderna.\n"
                "- Genocidio: Holocausto, genocidio de Ruanda, genocidio armenio.")

# Slide 6: Debate: Extensión de la Declaración de la ONU
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Debate: Extensión de la Declaración de la ONU"
content.text = ("- ¿Derechos Elementales vs. Derechos Extendidos?\n"
                "- ¿Debería la Declaración de la ONU enfocarse solo en los derechos básicos o incluir derechos más amplios como derechos económicos, sociales y culturales?")

# Slide 7: Derechos Humanos: Teoría vs. Realidad
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Derechos Humanos: Teoría vs. Realidad"
content.text = ("- Utopía de los Derechos Humanos: Idealismo versus realidad práctica.\n"
                "- Desafíos en la Implementación: Falta de voluntad política, conflictos culturales, intereses económicos.")

# Slide 8: Enfoque Regional: El Mundo Árabe
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Enfoque Regional: El Mundo Árabe"
content.text = ("- Situación de los Derechos Humanos en el Mundo Árabe: Casos específicos de derechos vulnerados.\n"
                "- Progreso y Retos: Reformas recientes, activismo local e internacional.")

# Slide 9: Análisis de Derechos Específicos
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Análisis de Derechos Específicos"
content.text = ("- Artículo 4: Prohibición de la Esclavitud: Historia y situación actual de la esclavitud, tráfico de personas, trabajo forzado, explotación de menores.\n"
                "- Artículo 5: Prohibición de la Tortura: Definición y ejemplos de tortura, trabajo forzado en la minería de cobalto, tráfico de personas.\n"
                "- Artículo 9: Derecho a no ser Detenido Arbitrariamente: Ejemplos de detenciones arbitrarias alrededor del mundo.\n"
                "- Artículo 12: Derecho a la Privacidad: Violaciones a la privacidad en la era digital.\n"
                "- Artículo 18: Libertad de Pensamiento, Conciencia y Religión: Ejemplos de violaciones a la libertad de pensamiento y religiosa.")

# Slide 10: Conclusión
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusión"
content.text = ("- Resumen de los Puntos Clave: Repaso de la evolución, interpretaciones, y violaciones de los derechos humanos.\n"
                "- Importancia de la Protección y Promoción de los Derechos Humanos: Necesidad de compromiso global para garantizar el respeto a los derechos humanos.")

# Slide 11: Preguntas y Debate
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Preguntas y Debate"
content.text = ("Abrir el espacio para preguntas del público y un debate sobre los temas discutidos.")

# Guardar la presentación
prs.save('presentacion_derechos_humanos.pptx')
